# scraping/files.py
import logging
import requests
from io import BytesIO

# Use PyPDF2 for basic PDF extraction, consider adding pdfminer.six as fallback
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
    logging.warning("PyPDF2 not installed. Basic PDF text extraction will not work.")

# Add imports for other file types as needed
try:
    import docx  # python-docx
except ImportError:
    docx = None
    logging.warning("python-docx not installed. DOCX extraction will not work.")

try:
    from pptx import Presentation  # python-pptx
except ImportError:
    Presentation = None
    logging.warning("python-pptx not installed. PPTX extraction will not work.")

from .core import create_session, make_request
from config import config  # Import the singleton instance

logger = logging.getLogger(__name__)


def fetch_file_content(username: str, password: str, file_url: str) -> bytes | None:
    """
    Fetches the binary content of a file using NTLM authentication.

    Args:
        username: User's GUC username.
        password: User's GUC password.
        file_url: The direct URL to the file.

    Returns:
        bytes: The file content on success.
        None: On failure (auth, network error, file not found).
    """
    session = create_session(username, password)
    logger.info(f"Attempting to fetch file content for {username} from {file_url}")

    # Use make_request which handles retries and basic errors
    # stream=True is important for potentially large files, though we read all content here
    # If streaming to client is needed, the API endpoint should handle response.iter_content
    response = make_request(
        session, file_url, method="GET", stream=False, timeout=(10, 60)
    )  # Longer read timeout for files

    if not response:
        logger.error(
            f"Failed to fetch file content from {file_url}. Make_request returned None."
        )
        # Specific error (401, timeout, etc.) logged by make_request
        return None

    # Check status code again (make_request might return None on 401, but double check)
    if response.status_code != 200:
        logger.error(
            f"Failed to fetch file content from {file_url}. Status code: {response.status_code}"
        )
        return None

    logger.info(
        f"Successfully fetched file content ({len(response.content)} bytes) from {file_url}"
    )
    return response.content


def extract_text_from_pdf(content: bytes) -> str:
    """Extracts text from PDF byte content using PyPDF2 (basic)."""
    if not PyPDF2:
        logger.error(
            "PyPDF2 library is required for PDF text extraction but not installed."
        )
        return "Error: PDF parsing library (PyPDF2) not installed."
    if not content:
        return ""

    text = ""
    try:
        pdf_file = BytesIO(content)
        reader = PyPDF2.PdfReader(pdf_file)
        num_pages = len(reader.pages)
        logger.info(
            f"Attempting to extract text from PDF with {num_pages} pages (using PyPDF2)."
        )
        page_texts = []
        for i in range(num_pages):
            try:
                page = reader.pages[i]
                page_text = page.extract_text()
                if page_text:
                    page_texts.append(page_text)
            except Exception as page_err:
                logger.warning(
                    f"PyPDF2: Error extracting text from PDF page {i+1}: {page_err}"
                )
                # Continue to next page

        text = "\n\n".join(page_texts)  # Join pages with double newline

        if not text.strip() and num_pages > 0:
            logger.warning("PyPDF2 extracted no text from the PDF.")
            # Consider adding pdfminer.six fallback here if needed and installed
            # text = _try_pdfminer_fallback(content)

    except PyPDF2.errors.PdfReadError as pdf_err:
        logger.error(f"PyPDF2 could not read PDF: {pdf_err}")
        text = "Error: Could not read PDF file (possibly corrupted or encrypted)."
    except Exception as e:
        logger.exception(f"Unexpected error during PyPDF2 PDF extraction: {e}")
        text = f"Error: Unexpected error during PDF processing: {e}"
    return text.strip()


# Placeholder for pdfminer fallback if you add it
# def _try_pdfminer_fallback(content: bytes) -> str:
#     try:
#         from pdfminer.high_level import extract_text as pdfminer_extract
#         logger.info("Using pdfminer.six fallback for PDF extraction.")
#         pdf_file = BytesIO(content)
#         return pdfminer_extract(pdf_file).strip()
#     except ImportError:
#         logger.error("pdfminer.six not installed, cannot use fallback.")
#         return "Error: PDF extraction failed (PyPDF2 found no text, pdfminer not installed)."
#     except Exception as pdfminer_err:
#         logger.error(f"pdfminer.six fallback failed: {pdfminer_err}")
#         return f"Error: PDF extraction failed (PyPDF2 and pdfminer errors: {pdfminer_err})."


def extract_text_from_docx(content: bytes) -> str:
    """Extracts text from DOCX byte content using python-docx."""
    if not docx:
        logger.error(
            "python-docx library is required for DOCX text extraction but not installed."
        )
        return "Error: DOCX parsing library (python-docx) not installed."
    if not content:
        return ""

    text = ""
    try:
        doc_file = BytesIO(content)
        document = docx.Document(doc_file)
        paragraphs = [para.text for para in document.paragraphs if para.text]
        text = "\n".join(paragraphs)
        logger.info(f"Extracted text from DOCX using python-docx.")
    except Exception as e:
        logger.exception(f"Error during python-docx DOCX extraction: {e}")
        text = f"Error: Failed to process DOCX file: {e}"
    return text.strip()


def extract_text_from_pptx(content: bytes) -> str:
    """Extracts text from PPTX byte content using python-pptx."""
    if not Presentation:
        logger.error(
            "python-pptx library is required for PPTX text extraction but not installed."
        )
        return "Error: PPTX parsing library (python-pptx) not installed."
    if not content:
        return ""

    text = ""
    try:
        ppt_file = BytesIO(content)
        prs = Presentation(ppt_file)
        slides_text = []
        logger.info(
            f"Attempting to extract text from PPTX with {len(prs.slides)} slides."
        )
        for i, slide in enumerate(prs.slides):
            slide_text_parts = []
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Add text from the shape, stripping extra whitespace
                        shape_text = shape.text.strip()
                        if shape_text:
                            slide_text_parts.append(shape_text)
            except Exception as shape_err:
                logger.warning(
                    f"Error processing shapes on PPTX slide {i+1}: {shape_err}"
                )
            if slide_text_parts:
                slides_text.append(
                    "\n".join(slide_text_parts)
                )  # Join text within a slide

        text = "\n\n".join(slides_text)  # Join text from different slides
        logger.info(f"Extracted text from PPTX using python-pptx.")
    except Exception as e:
        logger.exception(f"Error during python-pptx PPTX extraction: {e}")
        text = f"Error: Failed to process PPTX file: {e}"
    return text.strip()


def extract_text_from_file(content: bytes, filename: str) -> str:
    """
    Extracts text content from supported file types.

    Args:
        content: The binary content of the file.
        filename: The original filename, used to determine the file type.

    Returns:
        The extracted text as a string, or an error message if extraction fails
        or the type is unsupported.
    """
    if not content or not filename:
        logger.warning("extract_text_from_file called with empty content or filename.")
        return ""

    # Determine file extension
    parts = filename.lower().split("?")[0].split(".")
    extension = parts[-1] if len(parts) > 1 else ""

    logger.info(f"Attempting text extraction for file type: {extension}")

    if extension == "pdf":
        return extract_text_from_pdf(content)
    elif extension == "docx":
        return extract_text_from_docx(content)
    elif extension == "pptx":
        return extract_text_from_pptx(content)
    elif extension == "txt" or extension in [
        "html",
        "htm",
        "csv",
        "xml",
        "json",
        "js",
        "css",
    ]:  # Treat as plain text or structured text
        try:
            # Attempt decoding with common encodings
            try:
                return content.decode("utf-8").strip()
            except UnicodeDecodeError:
                try:
                    return content.decode("latin-1").strip()
                except UnicodeDecodeError:
                    logger.warning(
                        f"Could not decode text file '{filename}' with utf-8 or latin-1."
                    )
                    return "Error: Could not decode text file."
        except Exception as e:
            logger.error(f"Error reading text-based file '{filename}': {e}")
            return f"Error: Failed to read text file: {e}"
    else:
        logger.warning(f"Unsupported file type for text extraction: {extension}")
        return f"Unsupported file type for extraction: .{extension}"
